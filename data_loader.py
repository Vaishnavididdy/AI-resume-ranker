"""Enhanced Data loading — robust parsing with better error handling."""
from __future__ import annotations
import pandas as pd
import json
import os
import re
from typing import Any, Dict, Iterable, List, Optional
import traceback

# Global imports with fallback - use None if not available
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from docx import Document
except ImportError:
    Document = None


class FileParserError(Exception):
    """Custom exception for file parsing errors."""
    pass


def _get_document():
    """Lazy load Document to avoid startup issues."""
    global Document
    if Document is None:
        try:
            from docx import Document as DocxDocument
            Document = DocxDocument
            return Document
        except ImportError:
            raise FileParserError("python-docx not installed. Please install: pip install python-docx")
    return Document


def _get_pdf_reader():
    """Lazy load PdfReader to avoid startup issues."""
    global PdfReader
    if PdfReader is None:
        try:
            from pypdf import PdfReader as PR
            PdfReader = PR
            return PdfReader
        except ImportError:
            raise FileParserError("pypdf not installed. Please install: pip install pypdf")
    return PdfReader


def _get_presentation():
    """Lazy load Presentation to avoid startup issues."""
    global Presentation
    if Presentation is None:
        try:
            from pptx import Presentation as PR
            Presentation = PR
            return Presentation
        except ImportError:
            raise FileParserError("python-pptx not installed. Please install: pip install python-pptx")
    return Presentation


def load_job_description(path: str) -> Dict[str, Any]:
    """Enhanced job description loading with better error handling."""
    try:
        if not os.path.exists(path):
            raise FileParserError(f"File not found: {path}")
        
        ext = os.path.splitext(path)[1].lower()
        content = ""
        
        # Text-based formats
        if ext in [".txt", ".md"]:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        
        elif ext in [".json", ".jsonl"]:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                data = json.load(f)
                # Try to extract text content
                if isinstance(data, dict):
                    content = json.dumps(data, indent=2)
                elif isinstance(data, list):
                    content = json.dumps(data, indent=2)
                else:
                    content = str(data)
        
        elif ext in [".doc", ".docx"]:
            doc_module = _get_document()
            doc = doc_module(path)
            content = "\n".join(p.text for p in doc.paragraphs)
            # Also extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " ".join(cell.text for cell in row.cells)
                    content += "\n" + row_text
        
        elif ext == ".pdf":
            pdf_module = _get_pdf_reader()
            reader = pdf_module(path)
            content = "\n".join(page.extract_text() or "" for page in reader.pages)
            if not content.strip():
                raise FileParserError("No text extracted from PDF")
        
        elif ext == ".csv":
            df = pd.read_csv(path, encoding='utf-8', errors='ignore')
            content = df.to_string()
        
        elif ext in [".ppt", ".pptx"]:
            ppt_module = _get_presentation()
            prs = ppt_module(path)
            slides = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides.append(shape.text)
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            slides.append(paragraph.text)
            content = "\n".join(slides)
        
        elif ext in [".xlsx", ".xls"]:
            try:
                engine = 'openpyxl' if ext == '.xlsx' else 'xlrd'
                df = pd.read_excel(path, engine=engine)
                content = df.to_string()
            except Exception as e:
                raise FileParserError(f"Excel parsing failed: {e}")
        
        else:
            raise FileParserError(f"Unsupported file format: {ext}")
        
        if not content.strip():
            raise FileParserError(f"Empty or unreadable content from {path}")
        
        return {"raw_text": content, "fields": {}}
    
    except Exception as e:
        raise FileParserError(f"Failed to load job description: {str(e)}\n{traceback.format_exc()}")


def _iter_json_objects(text: str) -> Iterable[Dict[str, Any]]:
    """Enhanced JSON parsing with multiple format support."""
    text = text.strip()
    if not text:
        return
    
    # Try parsing as a JSON array
    if text.startswith("["):
        try:
            arr = json.loads(text)
            if isinstance(arr, list):
                for obj in arr:
                    if isinstance(obj, dict):
                        yield obj
                return
        except Exception:
            pass
    
    # Try parsing as JSONL (line by line)
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        
        try:
            obj = json.loads(line)
            if isinstance(obj, dict):
                yield obj
        except Exception:
            # Try to fix common issues
            try:
                # Try to parse multiple JSON objects on one line
                matches = re.findall(r'\{[^{}]*\}', line)
                for match in matches:
                    obj = json.loads(match)
                    if isinstance(obj, dict):
                        yield obj
            except Exception:
                continue


def load_candidates(path: str) -> List[Dict[str, Any]]:
    """Enhanced candidate loading with robust handling."""
    try:
        if not os.path.exists(path):
            raise FileParserError(f"File not found: {path}")
        
        ext = os.path.splitext(path)[1].lower()
        candidates = []
        
        # JSON/JSONL/TXT formats
        if ext in [".json", ".jsonl", ".txt"]:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            
            candidates = list(_iter_json_objects(text))
            
            # If no JSON objects found, try as plain text
            if not candidates and ext == ".txt":
                return [{"raw_text": text, "candidate_id": "text_only"}]
        
        # CSV format
        elif ext == ".csv":
            try:
                df = pd.read_csv(path, encoding='utf-8', errors='ignore')
                candidates = df.to_dict(orient="records")
            except UnicodeDecodeError:
                df = pd.read_csv(path, encoding='latin-1')
                candidates = df.to_dict(orient="records")
        
        # Excel formats
        elif ext in [".xlsx", ".xls"]:
            try:
                engine = 'openpyxl' if ext == '.xlsx' else 'xlrd'
                df = pd.read_excel(path, engine=engine)
                candidates = df.to_dict(orient="records")
            except Exception as e:
                raise FileParserError(f"Excel parsing failed: {e}")
        
        # DOC/DOCX
        elif ext in [".doc", ".docx"]:
            doc_module = _get_document()
            doc = doc_module(path)
            text = "\n".join(p.text for p in doc.paragraphs)
            candidates = [{"raw_text": text, "candidate_id": "docx_candidate"}]
        
        # PDF
        elif ext == ".pdf":
            pdf_module = _get_pdf_reader()
            reader = pdf_module(path)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            candidates = [{"raw_text": text, "candidate_id": "pdf_candidate"}]
        
        else:
            raise FileParserError(f"Unsupported file format: {ext}")
        
        # Ensure we have candidate IDs
        for i, cand in enumerate(candidates):
            if not cand.get("candidate_id") and not cand.get("id"):
                cand["candidate_id"] = f"candidate_{i+1:04d}"
        
        if not candidates:
            raise FileParserError("No candidates found in file")
        
        return candidates
    
    except Exception as e:
        raise FileParserError(f"Failed to load candidates: {str(e)}\n{traceback.format_exc()}")